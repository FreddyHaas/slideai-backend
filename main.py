from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import pandas as pd
from openai import OpenAI

from models import BarChartDataStructure, SelectedChartType

client = OpenAI()

CHART_CORE_MESSAGE = "China is the most important ice cream market in 2029"
DATA_PATH = "./salesdata_single_row_header.xlsx"
TEMPLATE_PATH = "./template.pptx"


# Query and mock responses

def _query_openai(message, response_model=None):
    if response_model is None:
        completion = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "user",
                    "content": message
                }
            ]
        )
        return completion.choices[0].message
    else:
        completion = client.beta.chat.completions.parse(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "user",
                    "content": message
                }
            ],
            response_format=response_model

        )
        return completion.choices[0].message.parsed


# Chart creator

def _create_stacked_bar_chart(pivot_table, diagram_title, chart_title):
    presentation = Presentation(TEMPLATE_PATH)
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])

    chart_data = CategoryChartData()
    categories_column = pivot_table.columns[0]
    chart_data.categories = pivot_table[categories_column].tolist()

    for column in pivot_table.columns[1:]:
        chart_data.add_series(str(column), pivot_table[column].tolist())

    diagram_placeholder = slide.placeholders[13]
    chart = diagram_placeholder.insert_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data
    ).chart

    # Customize the chart
    slide.shapes.title.text = chart_title
    chart.has_title = True
    chart.chart_title.text_frame.text = diagram_title
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM  # Bottom of the chart (default)

    # Save the presentation
    presentation.save("stacked_bar_chart.pptx")


# Prompts

def _create_chart_selection_prompt(table_headers, chart_message):
    return (f'Here are the columns of a table: {table_headers}. Which chart type is best suited to support the '
            f'message "{chart_message}". Please select from BAR_CHART and WATERFALL_CHART. '
            f'Respond exactly with that word.')


def _create_data_selection_prompt(table_headers, chart_message):
    return (f'Here are the columns of a table: {table_headers}. In order to create a stacked bar chart to support this '
            f'message: {chart_message}. Which column should be the category, subcategory and value? Please also provide '
            f'a short descriptive name for the values with a unit if possible (e.g. "living room size in square meter"')


# Main function
def create_chart(data_path, chart_core_message):
    df = pd.read_excel(data_path)

    df_headers = df.columns.tolist()

    chart_selection_prompt = _create_chart_selection_prompt(df_headers, chart_core_message)

    """selected_chart_type = _query_openai(
        message=chart_selection_prompt,
        response_model=SelectedChartType
    )"""
    # Mock
    selected_chart_type = SelectedChartType(chartType="BAR_CHART")

    if selected_chart_type.chartType == "BAR_CHART":

        data_selection_prompt = _create_data_selection_prompt(df_headers, chart_core_message)
        """bar_chart_data_structure = _query_openai(
            message=data_selection_prompt,
            response_model=BarChartDataStructure
        )"""
        # mock
        bar_chart_data_structure = BarChartDataStructure(
            category="Market",
            subcategory="Year",
            value="Ice cream sales",
            title="Ice cream sales in EUR"
        )

        pivot_df = df.pivot(
            index=bar_chart_data_structure.category,
            columns=bar_chart_data_structure.subcategory,
            values=bar_chart_data_structure.value
        )
        pivot_df = pivot_df.reset_index()  # Reset index for easier PowerPoint processing
        print(pivot_df)
        _create_stacked_bar_chart(pivot_df, bar_chart_data_structure.title, chart_core_message)


create_chart(DATA_PATH, CHART_CORE_MESSAGE)

