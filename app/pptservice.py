import datetime
import os

from openpyxl.reader.excel import load_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData, BubbleChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import pandas as pd
from openai import OpenAI
from pptx.util import Pt

from models import ClusteredBarOrColumnDataStructure, LineOrClusteredColumnChartDataStructure, SelectedChartType, \
    ChartType, \
    BubbleChartDataStructure, \
    BarOrColumnDataStructure, PieChartDataStructure, TablePivot

MOCK_AI_API_CALLS = False

LABEL_COLOR = RGBColor(89, 89, 89)
GRID_COLOR = RGBColor(217, 217, 217)
LABEL_FONT_SIZE = Pt(12)
LINE_WIDTH = Pt(0.2).emu

client = OpenAI()

current_dir = os.path.dirname(os.path.abspath(__file__))

CHART_CORE_MESSAGE = "China is the most important ice cream market in 2029"
DATA_THREE_COLUMNS = "./inputs/market_year_sales_single_row_header.xlsx"
DATA_TWO_COLUMNS = "./inputs/market_sales_single_row_header.xlsx"
DATA_PIE_CHART = "./inputs/pie_chart_sales_single_row_header.xlsx"
DATA_TWO_TIME_SERIES = "./inputs/two_time_series_single_row_header.xlsx"
DATA_BUBBLE_CHART = "./inputs/bubble_chart_single_row_header.xlsx"
TEMPLATE_PATH = os.path.join(current_dir, "template.pptx")


# ToDo: Sortierung bei diversen Charttypen (Bar, Column,...)
# ToDo: Verschiedene Daten einlesen

# Query and mock responses
def _query_openai(message, response_model=None):
    if response_model is None:
        completion = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": message
                }
            ],
            temperature=0
        )
        return completion.choices[0].message
    else:
        completion = client.beta.chat.completions.parse(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": message
                }
            ],
            temperature=0,
            response_format=response_model

        )
        return completion.choices[0].message.parsed


# Column chart creators
def _create_column_chart(slide, df, headers, chart_core_message, header_cell_formats):
    data_selection_prompt = _create_bar_or_column_chart_data_selection_prompt(
        table_headers=headers,
        chart_message=chart_core_message,
        chart_type="column chart",
        header_cell_formats=header_cell_formats)

    column_chart_data_structure = _query_openai(
        message=data_selection_prompt,
        response_model=BarOrColumnDataStructure
    ) if not MOCK_AI_API_CALLS else (
        BarOrColumnDataStructure(
            category="Market",
            value="Ice cream sales",
            title="Ice cream sales in EUR"
        )
    )

    df = df.groupby(column_chart_data_structure.category, as_index=False)[column_chart_data_structure.value].sum()

    # Prepare the data for the chart
    if column_chart_data_structure.category not in df.columns or \
            column_chart_data_structure.value not in df.columns:
        raise ValueError("Specified category or value columns not found in DataFrame.")

    # Chart creation
    chart_data = CategoryChartData()
    chart_data.categories = df[column_chart_data_structure.category].tolist()
    chart_data.add_series(
        column_chart_data_structure.value,
        df[column_chart_data_structure.value].tolist()
    )

    diagram_placeholder = slide.placeholders[13]
    chart = diagram_placeholder.insert_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data
    ).chart

    # Title and labels
    slide.shapes.title.text = chart_core_message
    chart.has_title = True
    chart.chart_title.text_frame.text = column_chart_data_structure.title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = LABEL_COLOR
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    # Gridlines
    value_axis = chart.value_axis
    major_gridlines = value_axis.major_gridlines
    line = major_gridlines.format.line
    line.fill.solid()
    line.fill.fore_color.rgb = GRID_COLOR
    line.width = LINE_WIDTH

    # Style the x-axis (category axis)
    category_axis = chart.category_axis

    # Format the x-axis line to match the major gridlines
    axis_line = category_axis.format.line
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = GRID_COLOR  # Same color as gridlines
    axis_line.width = LINE_WIDTH  # Same width as gridlines

    category_axis.tick_labels.font.size = LABEL_FONT_SIZE
    category_axis.tick_labels.font.color.rgb = LABEL_COLOR

    # Style the y-axis (value axis) but hide its line
    value_axis.tick_labels.font.size = LABEL_FONT_SIZE
    value_axis.tick_labels.font.color.rgb = LABEL_COLOR
    value_axis.format.line.fill.solid()
    value_axis.format.line.fill.background()  # Hide the axis line


def _create_clustered_column_chart(slide, df, headers, chart_core_message, header_cell_formats):

    pivot_prompt = _create_pivot_data_prompt(df=df, core_message=chart_core_message,
                                             chart_type="clustered column chart")

    pivot_response = _query_openai(
        message=pivot_prompt,
        response_model=TablePivot
    ) if not MOCK_AI_API_CALLS else (
        TablePivot(
            needsPivoting=False,
            index="null",
            columns="null",
            values="null"
        )
    )

    print(pivot_response.data_structure_analysis)
    print(pivot_response.explain_pivot)

    if pivot_response.long_format is True:
        df = df.pivot(index=pivot_response.index,
                      columns=pivot_response.columns,
                      values=pivot_response.values
                      )
        headers = df.columns.tolist()
        values_cell_formats = header_cell_formats[pivot_response.values]
        header_cell_formats = {}
        for header in headers:
            header_cell_formats[header] = f"{pivot_response.values} in format ${values_cell_formats}"

    data_selection_prompt = _create_line_or_clustered_column_chart_data_selection_prompt(headers,
                                                                                         chart_core_message,
                                                                                         "clustered column chart",
                                                                                         header_cell_formats)

    column_chart_data_structure = _query_openai(
        message=data_selection_prompt,
        response_model=LineOrClusteredColumnChartDataStructure
    ) if not MOCK_AI_API_CALLS else (
        LineOrClusteredColumnChartDataStructure(
            category="Year",
            series=["USA", "China"],
            title="Some title"
        )
    )

    # Chart creation
    chart_data = CategoryChartData()
    chart_data.categories = df[
        column_chart_data_structure.category].tolist() if not pivot_response.long_format else df.index.tolist()

    # Convert all column headers to strings to avoid errors when matching with selected columns from openai prompt
    df.columns = df.columns.astype(str)

    for column in column_chart_data_structure.series:
        chart_data.add_series(column, df[column].tolist())

    diagram_placeholder = slide.placeholders[13]
    chart = diagram_placeholder.insert_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data
    ).chart

    # Title and labels
    slide.shapes.title.text = chart_core_message
    chart.has_title = True
    chart.chart_title.text_frame.text = column_chart_data_structure.title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = LABEL_COLOR
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    # Chart legend
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM  # Bottom of the chart (default)
    chart.legend.font.size = LABEL_FONT_SIZE

    # Gridlines
    value_axis = chart.value_axis
    major_gridlines = value_axis.major_gridlines
    line = major_gridlines.format.line
    line.fill.solid()
    line.fill.fore_color.rgb = GRID_COLOR
    line.width = LINE_WIDTH

    # Style the x-axis (category axis)
    category_axis = chart.category_axis

    # Format the x-axis line to match the major gridlines
    axis_line = category_axis.format.line
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = GRID_COLOR  # Same color as gridlines
    axis_line.width = LINE_WIDTH  # Same width as gridlines

    category_axis.tick_labels.font.size = LABEL_FONT_SIZE
    category_axis.tick_labels.font.color.rgb = LABEL_COLOR

    # Style the y-axis (value axis) but hide its line
    value_axis.tick_labels.font.size = LABEL_FONT_SIZE
    value_axis.tick_labels.font.color.rgb = LABEL_COLOR
    value_axis.format.line.fill.solid()
    value_axis.format.line.fill.background()  # Hide the axis line

    # Adjust overlap to add spacing between columns in the same cluster
    chart.plots[0].overlap = -25


def _create_stacked_column_chart(slide, df, headers, chart_core_message, header_cell_formats):
    pivot_prompt = _create_pivot_data_prompt(df=df, core_message=chart_core_message,
                                             chart_type="clustered column chart")

    pivot_response = _query_openai(
        message=pivot_prompt,
        response_model=TablePivot
    ) if not MOCK_AI_API_CALLS else (
        TablePivot(
            needsPivoting=False,
            index="null",
            columns="null",
            values="null"
        )
    )

    print(pivot_response.data_structure_analysis)
    print(pivot_response.explain_pivot)

    if pivot_response.long_format is True:
        df = df.pivot(index=pivot_response.index,
                      columns=pivot_response.columns,
                      values=pivot_response.values
                      )
        headers = df.columns.tolist()
        values_cell_formats = header_cell_formats[pivot_response.values]
        header_cell_formats = {}
        for header in headers:
            header_cell_formats[header] = f"{pivot_response.values} in format ${values_cell_formats}"

    data_selection_prompt = _create_line_or_clustered_column_chart_data_selection_prompt(headers,
                                                                                         chart_core_message,
                                                                                         "clustered column chart",
                                                                                         header_cell_formats)

    column_chart_data_structure = _query_openai(
        message=data_selection_prompt,
        response_model=LineOrClusteredColumnChartDataStructure
    ) if not MOCK_AI_API_CALLS else (
        LineOrClusteredColumnChartDataStructure(
            category="Year",
            series=["USA", "China"],
            title="Some title"
        )
    )

    # Chart creation
    chart_data = CategoryChartData()
    chart_data.categories = df[
        column_chart_data_structure.category].tolist() if not pivot_response.long_format else df.index.tolist()

    # Convert all column headers to strings to avoid errors when matching with selected columns from openai prompt
    df.columns = df.columns.astype(str)

    for column in df.columns[1:]:
        chart_data.add_series(str(column), df[column].tolist())

    diagram_placeholder = slide.placeholders[13]
    chart = diagram_placeholder.insert_chart(
        XL_CHART_TYPE.COLUMN_STACKED, chart_data  # Set chart type to stacked column
    ).chart

    # Title and labels
    slide.shapes.title.text = chart_core_message
    chart.has_title = True
    chart.chart_title.text_frame.text = column_chart_data_structure.title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = LABEL_COLOR
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    # Chart legend
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM  # Bottom of the chart (default)
    chart.legend.font.size = LABEL_FONT_SIZE

    # Gridlines
    value_axis = chart.value_axis
    major_gridlines = value_axis.major_gridlines
    line = major_gridlines.format.line
    line.fill.solid()
    line.fill.fore_color.rgb = GRID_COLOR
    line.width = LINE_WIDTH

    # Style the x-axis (category axis)
    category_axis = chart.category_axis

    # Format the x-axis line to match the major gridlines
    axis_line = category_axis.format.line
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = GRID_COLOR  # Same color as gridlines
    axis_line.width = LINE_WIDTH  # Same width as gridlines

    category_axis.tick_labels.font.size = LABEL_FONT_SIZE
    category_axis.tick_labels.font.color.rgb = LABEL_COLOR

    # Style the y-axis (value axis) but hide its line
    value_axis.tick_labels.font.size = LABEL_FONT_SIZE
    value_axis.tick_labels.font.color.rgb = LABEL_COLOR
    value_axis.format.line.fill.solid()
    value_axis.format.line.fill.background()  # Hide the axis line


def _create_100_percent_stacked_column_chart(slide, df, headers, chart_core_message, header_cell_formats):
    pivot_prompt = _create_pivot_data_prompt(df=df, core_message=chart_core_message,
                                             chart_type="100% stacked column chart")

    pivot_response = _query_openai(
        message=pivot_prompt,
        response_model=TablePivot
    ) if not MOCK_AI_API_CALLS else (
        TablePivot(
            needsPivoting=False,
            index="null",
            columns="null",
            values="null"
        )
    )

    print(pivot_response.data_structure_analysis)
    print(pivot_response.explain_pivot)

    if pivot_response.long_format is True:
        df = df.pivot(index=pivot_response.index,
                      columns=pivot_response.columns,
                      values=pivot_response.values
                      )
        headers = df.columns.tolist()
        values_cell_formats = header_cell_formats[pivot_response.values]
        header_cell_formats = {}
        for header in headers:
            header_cell_formats[header] = f"{pivot_response.values} in format ${values_cell_formats}"

    data_selection_prompt = _create_line_or_clustered_column_chart_data_selection_prompt(headers,
                                                                                         chart_core_message,
                                                                                         "100% stacked column chart",
                                                                                         header_cell_formats)

    column_chart_data_structure = _query_openai(
        message=data_selection_prompt,
        response_model=LineOrClusteredColumnChartDataStructure
    ) if not MOCK_AI_API_CALLS else (
        LineOrClusteredColumnChartDataStructure(
            category="Year",
            series=["USA", "China"],
            title="Some title"
        )
    )

    df.reset_index()

    # Normalize the pivoted data to get 100% stacked values (percentage)
    pivot_df_percentage = df.copy()
    pivot_df_percentage[column_chart_data_structure.series] = df[column_chart_data_structure.series].div(df[column_chart_data_structure.series].sum(axis=1), axis=0) * 100

    # Chart creation
    chart_data = CategoryChartData()
    chart_data.categories = pivot_df_percentage[column_chart_data_structure.category].tolist() if not pivot_response.long_format else df.index.tolist()

    # Convert all column headers to strings to avoid errors when matching with selected columns from openai prompt
    df.columns = df.columns.astype(str)

    for column in column_chart_data_structure.series:
        chart_data.add_series(column, pivot_df_percentage[column].tolist())

    diagram_placeholder = slide.placeholders[13]
    chart = diagram_placeholder.insert_chart(
        XL_CHART_TYPE.COLUMN_STACKED, chart_data  # Set chart type to stacked column
    ).chart

    # Title and labels
    slide.shapes.title.text = chart_core_message
    chart.has_title = True
    chart.chart_title.text_frame.text = column_chart_data_structure.title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = LABEL_COLOR
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    # Chart legend
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM  # Bottom of the chart (default)
    chart.legend.font.size = LABEL_FONT_SIZE

    # Gridlines
    value_axis = chart.value_axis
    major_gridlines = value_axis.major_gridlines
    line = major_gridlines.format.line
    line.fill.solid()
    line.fill.fore_color.rgb = GRID_COLOR
    line.width = LINE_WIDTH

    # Style the x-axis (category axis)
    category_axis = chart.category_axis

    # Format the x-axis line to match the major gridlines
    axis_line = category_axis.format.line
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = GRID_COLOR  # Same color as gridlines
    axis_line.width = LINE_WIDTH  # Same width as gridlines

    category_axis.tick_labels.font.size = LABEL_FONT_SIZE
    category_axis.tick_labels.font.color.rgb = LABEL_COLOR

    # Style the y-axis (value axis) but hide its line
    value_axis.tick_labels.font.size = LABEL_FONT_SIZE
    value_axis.tick_labels.font.color.rgb = LABEL_COLOR
    value_axis.format.line.fill.solid()
    value_axis.format.line.fill.background()  # Hide the axis line

    # Set the y-axis to range from 0 to 100 (for 100% stacked chart)
    value_axis.maximum_scale = 100
    value_axis.minimum_scale = 0
    value_axis.tick_labels.show_percentage = True
    value_axis.tick_labels.number_format = "0\%"  # This ensures the labels show as percentages


# Bar chart creators
def _create_bar_chart(slide, df, headers, chart_core_message, header_cell_formats):
    data_selection_prompt = _create_bar_or_column_chart_data_selection_prompt(
        table_headers=headers,
        chart_message=chart_core_message,
        chart_type="bar chart",
        header_cell_formats=header_cell_formats)

    bar_chart_data_structure = _query_openai(
        message=data_selection_prompt,
        response_model=BarOrColumnDataStructure
    ) if not MOCK_AI_API_CALLS else (
        BarOrColumnDataStructure(
            category="Market",
            value="Ice cream sales",
            title="Ice cream sales by market"
        )
    )

    # Prepare the data for the chart
    if bar_chart_data_structure.category not in df.columns or \
            bar_chart_data_structure.value not in df.columns:
        raise ValueError("Specified category or value columns not found in DataFrame.")

    # Convert all column headers to strings to avoid errors when matching with selected columns from openai prompt
    df.columns = df.columns.astype(str)

    # Chart creation
    chart_data = CategoryChartData()
    chart_data.categories = df[bar_chart_data_structure.category].tolist()
    chart_data.add_series(
        bar_chart_data_structure.value,
        df[bar_chart_data_structure.value].tolist()
    )

    diagram_placeholder = slide.placeholders[13]
    chart = diagram_placeholder.insert_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, chart_data
    ).chart

    # Title and labels
    slide.shapes.title.text = chart_core_message
    chart.has_title = True
    chart.chart_title.text_frame.text = bar_chart_data_structure.title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = LABEL_COLOR
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    # Gridlines
    value_axis = chart.value_axis
    major_gridlines = value_axis.major_gridlines
    line = major_gridlines.format.line
    line.fill.solid()
    line.fill.fore_color.rgb = GRID_COLOR
    line.width = LINE_WIDTH

    # Style the y-axis (category axis)
    category_axis = chart.category_axis

    # Format the y-axis line to match the major gridlines
    axis_line = category_axis.format.line
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = GRID_COLOR
    axis_line.width = LINE_WIDTH

    category_axis.tick_labels.font.size = LABEL_FONT_SIZE
    category_axis.tick_labels.font.color.rgb = LABEL_COLOR

    # Style the x-axis (value axis)
    value_axis.tick_labels.font.size = LABEL_FONT_SIZE
    value_axis.tick_labels.font.color.rgb = LABEL_COLOR
    value_axis.format.line.fill.solid()
    value_axis.format.line.fill.background()  # Hide the axis line


def _create_clustered_bar_chart(slide, df, headers, chart_core_message, header_cell_formats):
    data_selection_prompt = _create_clustered_bar_or_column_chart_data_selection_prompt(
        table_headers=headers,
        chart_message=chart_core_message,
        chart_type="clustered column chart",
        header_cell_formats=header_cell_formats)
    bar_chart_data_structure = _query_openai(
        message=data_selection_prompt,
        response_model=ClusteredBarOrColumnDataStructure
    ) if not MOCK_AI_API_CALLS else (
        ClusteredBarOrColumnDataStructure(
            category="Market",
            subcategory="Year",
            value="Ice cream sales",
            title="Ice cream sales by market"
        )
    )

    pivot_df = df.pivot(
        index=bar_chart_data_structure.category,
        columns=bar_chart_data_structure.subcategory,
        values=bar_chart_data_structure.value
    )
    pivot_df = pivot_df.reset_index()  # Reset index for easier PowerPoint processing
    pivot_df = pivot_df.sort_values(by=pivot_df.columns[-1], ascending=True)

    # Chart creation
    chart_data = CategoryChartData()
    categories_column = pivot_df.columns[0]
    chart_data.categories = pivot_df[categories_column].tolist()

    for column in pivot_df.columns[1:]:
        chart_data.add_series(str(column), pivot_df[column].tolist())

    diagram_placeholder = slide.placeholders[13]
    chart = diagram_placeholder.insert_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, chart_data
    ).chart

    # Title and labels
    slide.shapes.title.text = chart_core_message
    chart.has_title = True
    chart.chart_title.text_frame.text = bar_chart_data_structure.title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = LABEL_COLOR
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    # Chart legend
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = LABEL_FONT_SIZE

    # Gridlines
    value_axis = chart.value_axis
    major_gridlines = value_axis.major_gridlines
    line = major_gridlines.format.line
    line.fill.solid()
    line.fill.fore_color.rgb = GRID_COLOR
    line.width = LINE_WIDTH

    # Style the y-axis (category axis)
    category_axis = chart.category_axis

    # Format the y-axis line to match the major gridlines
    axis_line = category_axis.format.line
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = GRID_COLOR
    axis_line.width = LINE_WIDTH

    category_axis.tick_labels.font.size = LABEL_FONT_SIZE
    category_axis.tick_labels.font.color.rgb = LABEL_COLOR

    # Style the x-axis (value axis)
    value_axis.tick_labels.font.size = LABEL_FONT_SIZE
    value_axis.tick_labels.font.color.rgb = LABEL_COLOR
    value_axis.format.line.fill.solid()
    value_axis.format.line.fill.background()  # Hide the axis line


def _create_stacked_bar_chart(slide, df, headers, chart_core_message, header_cell_formats):
    data_selection_prompt = _create_clustered_bar_or_column_chart_data_selection_prompt(
        table_headers=headers,
        chart_message=chart_core_message,
        chart_type="clustered column chart",
        header_cell_formats=header_cell_formats)
    bar_chart_data_structure = _query_openai(
        message=data_selection_prompt,
        response_model=ClusteredBarOrColumnDataStructure
    ) if not MOCK_AI_API_CALLS else (
        ClusteredBarOrColumnDataStructure(
            category="Market",
            subcategory="Year",
            value="Ice cream sales",
            title="Ice cream sales by market"
        )
    )

    # Pivot the dataframe to structure data for stacking
    pivot_df = df.pivot(
        index=bar_chart_data_structure.category,
        columns=bar_chart_data_structure.subcategory,
        values=bar_chart_data_structure.value
    )
    pivot_df = pivot_df.reset_index()  # Reset index for easier PowerPoint processing

    # Prepare the data for the chart
    chart_data = CategoryChartData()
    categories_column = pivot_df.columns[0]
    chart_data.categories = pivot_df[categories_column].tolist()

    # Add each subcategory as a series for stacking
    for column in pivot_df.columns[1:]:
        chart_data.add_series(str(column), pivot_df[column].tolist())

    # Insert a stacked bar chart
    diagram_placeholder = slide.placeholders[13]
    chart = diagram_placeholder.insert_chart(
        XL_CHART_TYPE.BAR_STACKED, chart_data
    ).chart

    # Title and labels
    slide.shapes.title.text = chart_core_message
    chart.has_title = True
    chart.chart_title.text_frame.text = bar_chart_data_structure.title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = LABEL_COLOR
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    # Chart legend
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = LABEL_FONT_SIZE

    # Gridlines
    value_axis = chart.value_axis
    major_gridlines = value_axis.major_gridlines
    line = major_gridlines.format.line
    line.fill.solid()
    line.fill.fore_color.rgb = GRID_COLOR
    line.width = LINE_WIDTH

    # Style the y-axis (category axis)
    category_axis = chart.category_axis
    axis_line = category_axis.format.line
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = GRID_COLOR
    axis_line.width = LINE_WIDTH

    category_axis.tick_labels.font.size = LABEL_FONT_SIZE
    category_axis.tick_labels.font.color.rgb = LABEL_COLOR

    # Style the x-axis (value axis)
    value_axis.tick_labels.font.size = LABEL_FONT_SIZE
    value_axis.tick_labels.font.color.rgb = LABEL_COLOR
    value_axis.format.line.fill.solid()
    value_axis.format.line.fill.background()  # Hide the axis line

    # Adjust gap width for stacking aesthetics
    chart.plots[0].gap_width = 50  # Adjust as needed for aesthetics


def _create_100_percent_stacked_bar_chart(slide, df, headers, chart_core_message, header_cell_formats):
    data_selection_prompt = _create_clustered_bar_or_column_chart_data_selection_prompt(
        table_headers=headers,
        chart_message=chart_core_message,
        chart_type="clustered column chart",
        header_cell_formats=header_cell_formats)
    bar_chart_data_structure = _query_openai(
        message=data_selection_prompt,
        response_model=ClusteredBarOrColumnDataStructure
    ) if not MOCK_AI_API_CALLS else (
        ClusteredBarOrColumnDataStructure(
            category="Market",
            subcategory="Year",
            value="Ice cream sales",
            title="Ice cream sales by market"
        )
    )

    # Pivot the dataframe to structure data for stacking
    pivot_df = df.pivot(
        index=bar_chart_data_structure.category,
        columns=bar_chart_data_structure.subcategory,
        values=bar_chart_data_structure.value
    )
    pivot_df = pivot_df.reset_index()  # Reset index for easier PowerPoint processing
    pivot_df = pivot_df.sort_values(by=pivot_df.columns[-1], ascending=True)

    # Normalize the values to percentages for 100% stacked bar chart
    pivot_df_percentage = pivot_df.copy()
    pivot_df_percentage[pivot_df.columns[1:]] = pivot_df[pivot_df.columns[1:]].div(
        pivot_df[pivot_df.columns[1:]].sum(axis=1), axis=0
    ) * 100  # Convert to percentage

    # Prepare the data for the chart
    chart_data = CategoryChartData()
    categories_column = pivot_df_percentage.columns[0]
    chart_data.categories = pivot_df_percentage[categories_column].tolist()

    # Add each subcategory as a series for stacking (now with percentage values)
    for column in pivot_df_percentage.columns[1:]:
        chart_data.add_series(str(column), pivot_df_percentage[column].tolist())

    # Insert a 100% stacked bar chart
    diagram_placeholder = slide.placeholders[13]
    chart = diagram_placeholder.insert_chart(
        XL_CHART_TYPE.BAR_STACKED_100, chart_data
    ).chart

    # Title and labels
    slide.shapes.title.text = chart_core_message
    chart.has_title = True
    chart.chart_title.text_frame.text = bar_chart_data_structure.title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = LABEL_COLOR
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    # Chart legend
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = LABEL_FONT_SIZE

    # Gridlines
    value_axis = chart.value_axis
    major_gridlines = value_axis.major_gridlines
    line = major_gridlines.format.line
    line.fill.solid()
    line.fill.fore_color.rgb = GRID_COLOR
    line.width = LINE_WIDTH

    # Style the y-axis (category axis)
    category_axis = chart.category_axis
    axis_line = category_axis.format.line
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = GRID_COLOR
    axis_line.width = LINE_WIDTH

    category_axis.tick_labels.font.size = LABEL_FONT_SIZE
    category_axis.tick_labels.font.color.rgb = LABEL_COLOR

    # Style the x-axis (value axis)
    value_axis.tick_labels.font.size = LABEL_FONT_SIZE
    value_axis.tick_labels.font.color.rgb = LABEL_COLOR
    value_axis.format.line.fill.solid()
    value_axis.format.line.fill.background()  # Hide the axis line

    # Adjust gap width for stacking aesthetics
    chart.plots[0].gap_width = 50  # Adjust as needed for aesthetics


# Pie chart creators
# ToDo: Label format
def _create_pie_chart(slide, df, headers, chart_core_message, header_cell_formats):
    data_selection_prompt = _create_pie_chart_data_selection_prompt(
        table_headers=headers,
        chart_message=chart_core_message,
        chart_type="pie chart",
        header_cell_formats=header_cell_formats)
    chart_data_structure = _query_openai(
        message=data_selection_prompt,
        response_model=PieChartDataStructure
    ) if not MOCK_AI_API_CALLS else (
        PieChartDataStructure(
            category_column="Market",
            value_column="Ice cream sales",
            title="Ice cream sales by market"
        )
    )

    # Calculate percentages from the values
    total = df[chart_data_structure.value_column].sum()
    df_percentages = df.copy()
    df_percentages[chart_data_structure.value_column] = (df[chart_data_structure.value_column] / total)

    # Convert all column headers to strings to avoid errors when matching with selected columns from openai prompt
    df.columns = df.columns.astype(str)

    # Create a CategoryChartData object
    chart_data = CategoryChartData()
    chart_data.categories = df_percentages[chart_data_structure.category_column].tolist()
    chart_data.add_series('Percentage', df_percentages[chart_data_structure.value_column].tolist())

    # Create the pie chart
    diagram_placeholder = slide.placeholders[13]
    chart = diagram_placeholder.insert_chart(XL_CHART_TYPE.PIE, chart_data).chart

    # Title and labels
    slide.shapes.title.text = chart_core_message
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_data_structure.title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = LABEL_COLOR
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    # Chart legend
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(10)  # You can customize this value

    # Adding data labels to the pie chart
    plot = chart.plots[0]  # Pie chart typically has one plot
    plot.has_data_labels = True  # Enable data labels
    data_labels = plot.data_labels
    data_labels.show_value = True  # Show values
    data_labels.number_format = "0%"  # Customize the format, e.g., "Percentage" for percent values
    data_labels.font.size = Pt(18)  # Customize the font size of the labels

    return chart


def _create_doughnut_chart(slide, df, headers, chart_core_message, header_cell_formats):
    data_selection_prompt = _create_pie_chart_data_selection_prompt(
        table_headers=headers,
        chart_message=chart_core_message,
        chart_type="doughnut chart",
        header_cell_formats=header_cell_formats)
    chart_data_structure = _query_openai(
        message=data_selection_prompt,
        response_model=PieChartDataStructure
    ) if not MOCK_AI_API_CALLS else (
        PieChartDataStructure(
            category_column="Market",
            value_column="Ice cream sales",
            title="Ice cream sales by market"
        )
    )

    # Calculate percentages from the values
    total = df[chart_data_structure.value_column].sum()
    df_percentages = df.copy()
    df_percentages[chart_data_structure.value_column] = (df[chart_data_structure.value_column] / total)

    # Convert all column headers to strings to avoid errors when matching with selected columns from openai prompt
    df.columns = df.columns.astype(str)

    # Create a CategoryChartData object
    chart_data = CategoryChartData()
    chart_data.categories = df_percentages[chart_data_structure.category_column].tolist()
    chart_data.add_series('Percentage', df_percentages[chart_data_structure.value_column].tolist())

    # Create the pie chart
    diagram_placeholder = slide.placeholders[13]
    chart = diagram_placeholder.insert_chart(XL_CHART_TYPE.DOUGHNUT, chart_data).chart

    # Title and labels
    slide.shapes.title.text = chart_core_message
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_data_structure.title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = LABEL_COLOR
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    # Chart legend
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(10)  # You can customize this value

    # Adding data labels to the pie chart
    plot = chart.plots[0]  # Pie chart typically has one plot
    plot.has_data_labels = True  # Enable data labels
    data_labels = plot.data_labels
    data_labels.show_value = True  # Show values
    data_labels.number_format = "0%"  # Customize the format, e.g., "Percentage" for percent values
    data_labels.font.size = Pt(18)  # Customize the font size of the labels

    return chart


# Time series data
def _create_line_chart(slide, df, headers, chart_core_message, header_cell_formats):
    more_than_two_columns = len(headers) > 2
    needs_pivoting = False

    if more_than_two_columns:
        pivot_prompt = _create_pivot_data_prompt(df=df, core_message=chart_core_message, chart_type="line chart")

        pivot_response = _query_openai(
            message=pivot_prompt,
            response_model=TablePivot
        ) if not MOCK_AI_API_CALLS else (
            TablePivot(
                data_structure_analysis="",
                long_format=False,
                explain_pivot="",
                index="",
                columns="",
                values=""
            )
        )

        print(pivot_response.data_structure_analysis)
        print(pivot_response.explain_pivot)

        needs_pivoting = pivot_response.long_format

        if needs_pivoting is True:
            df = df.pivot(index=pivot_response.index,
                          columns=pivot_response.columns,
                          values=pivot_response.values
                          )
            headers = df.columns.tolist()
            values_cell_formats = header_cell_formats[pivot_response.values]
            header_cell_formats = {}
            for header in headers:
                header_cell_formats[header] = f"{pivot_response.values} in format ${values_cell_formats}"

    data_selection_prompt = _create_line_or_clustered_column_chart_data_selection_prompt(headers,
                                                                                         chart_core_message,
                                                                                         "line chart",
                                                                                         header_cell_formats)

    line_chart_data_structure = _query_openai(
        message=data_selection_prompt,
        response_model=LineOrClusteredColumnChartDataStructure
    ) if not MOCK_AI_API_CALLS else (
        LineOrClusteredColumnChartDataStructure(
            category="Zeilenbeschriftungen",
            series=["USA", "China"],
            title="Some title"
        )
    )

    if more_than_two_columns is False or needs_pivoting is False:
        df = df.groupby(line_chart_data_structure.category, as_index=False)[line_chart_data_structure.series].sum()

    # Chart creation
    chart_data = CategoryChartData()
    chart_data.categories = df[line_chart_data_structure.category].tolist() if not needs_pivoting else df.index.tolist()

    # Convert all column headers to strings to avoid errors when matching with selected columns from openai prompt
    df.columns = df.columns.astype(str)

    for column in line_chart_data_structure.series:
        chart_data.add_series(column, df[column].tolist())

    diagram_placeholder = slide.placeholders[13]
    chart = diagram_placeholder.insert_chart(
        XL_CHART_TYPE.LINE, chart_data
    ).chart

    # Title and labels
    slide.shapes.title.text = chart_core_message
    chart.has_title = True
    chart.chart_title.text_frame.text = line_chart_data_structure.title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = LABEL_COLOR
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    # Chart legend
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = LABEL_FONT_SIZE

    # Gridlines
    value_axis = chart.value_axis
    major_gridlines = value_axis.major_gridlines
    line = major_gridlines.format.line
    line.fill.solid()
    line.fill.fore_color.rgb = GRID_COLOR
    line.width = LINE_WIDTH

    # Style the y-axis (category axis)
    category_axis = chart.category_axis

    # Format the y-axis line to match the major gridlines
    axis_line = category_axis.format.line
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = GRID_COLOR
    axis_line.width = LINE_WIDTH

    category_axis.tick_labels.font.size = LABEL_FONT_SIZE
    category_axis.tick_labels.font.color.rgb = LABEL_COLOR

    # Style the x-axis (value axis)
    value_axis.tick_labels.font.size = LABEL_FONT_SIZE
    value_axis.tick_labels.font.color.rgb = LABEL_COLOR
    value_axis.format.line.fill.solid()
    value_axis.format.line.fill.background()  # Hide the axis line


def _create_stacked_area_chart(slide, df, headers, chart_core_message, header_cell_formats):
    data_selection_prompt = _create_line_or_clustered_column_chart_data_selection_prompt(headers, chart_core_message,
                                                                                         "stacked area chart",
                                                                                         header_cell_formats)
    line_chart_data_structure = _query_openai(
        message=data_selection_prompt,
        response_model=LineOrClusteredColumnChartDataStructure
    ) if not MOCK_AI_API_CALLS else (
        LineOrClusteredColumnChartDataStructure(
            category="Year",
            series=["USA", "China"],
            title="Some title"
        )
    )

    # Chart creation
    chart_data = CategoryChartData()
    chart_data.categories = df[line_chart_data_structure.category].tolist()

    # Convert all column headers to strings to avoid errors when matching with selected columns from openai prompt
    df.columns = df.columns.astype(str)

    for column in line_chart_data_structure.series:
        chart_data.add_series(column, df[column].tolist())

    diagram_placeholder = slide.placeholders[13]
    chart = diagram_placeholder.insert_chart(
        XL_CHART_TYPE.AREA_STACKED, chart_data
    ).chart

    # Title and labels
    slide.shapes.title.text = chart_core_message
    chart.has_title = True
    chart.chart_title.text_frame.text = line_chart_data_structure.title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = LABEL_COLOR
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    # Chart legend
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = LABEL_FONT_SIZE

    # Gridlines
    value_axis = chart.value_axis
    major_gridlines = value_axis.major_gridlines
    line = major_gridlines.format.line
    line.fill.solid()
    line.fill.fore_color.rgb = GRID_COLOR
    line.width = LINE_WIDTH

    # Style the y-axis (category axis)
    category_axis = chart.category_axis

    # Format the y-axis line to match the major gridlines
    axis_line = category_axis.format.line
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = GRID_COLOR
    axis_line.width = LINE_WIDTH

    category_axis.tick_labels.font.size = LABEL_FONT_SIZE
    category_axis.tick_labels.font.color.rgb = LABEL_COLOR

    # Style the x-axis (value axis)
    value_axis.tick_labels.font.size = LABEL_FONT_SIZE
    value_axis.tick_labels.font.color.rgb = LABEL_COLOR
    value_axis.format.line.fill.solid()
    value_axis.format.line.fill.background()  # Hide the axis line


# ToDo: Title anpassen
# ToDo: Styling
def _create_100_percent_stacked_area_chart(slide, df, headers, chart_core_message, header_cell_formats):
    data_selection_prompt = _create_line_or_clustered_column_chart_data_selection_prompt(headers, chart_core_message,
                                                                                         "stacked area chart",
                                                                                         header_cell_formats)
    line_chart_data_structure = _query_openai(
        message=data_selection_prompt,
        response_model=LineOrClusteredColumnChartDataStructure
    ) if not MOCK_AI_API_CALLS else (
        LineOrClusteredColumnChartDataStructure(
            category="Year",
            series=["USA", "China"],
            title="Some title"
        )
    )

    # Normalize the data to make sure each row sums to 100%
    normalized_df = df.copy()
    rows_sum = df.drop(columns=[line_chart_data_structure.category]).sum(axis=1)
    for column in line_chart_data_structure.series:
        normalized_df[column] = df[column] / rows_sum * 100

    # Chart creation
    chart_data = CategoryChartData()
    chart_data.categories = normalized_df[line_chart_data_structure.category].tolist()

    # Convert all column headers to strings to avoid errors when matching with selected columns from openai prompt
    df.columns = df.columns.astype(str)

    # Adding series to the chart
    for column in line_chart_data_structure.series:
        chart_data.add_series(column, normalized_df[column].tolist())

    diagram_placeholder = slide.placeholders[13]
    chart = diagram_placeholder.insert_chart(
        XL_CHART_TYPE.AREA_STACKED, chart_data
    ).chart

    # Title and labels
    slide.shapes.title.text = chart_core_message
    chart.has_title = True
    chart.chart_title.text_frame.text = line_chart_data_structure.title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = LABEL_COLOR
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    # Chart legend
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = LABEL_FONT_SIZE

    # Style the x-axis (category axis)
    category_axis = chart.category_axis
    axis_line = category_axis.format.line
    axis_line.fill.solid()
    axis_line.width = LINE_WIDTH

    category_axis.tick_labels.font.size = LABEL_FONT_SIZE
    category_axis.tick_labels.font.color.rgb = LABEL_COLOR

    # Style the y-axis (value axis)
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = LABEL_FONT_SIZE
    value_axis.format.line.fill.solid()
    value_axis.format.line.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White (hidden axis line)

    # Set the y-axis to range from 0 to 100 (for 100% stacked chart)
    value_axis.maximum_scale = 100
    value_axis.minimum_scale = 0
    value_axis.tick_labels.show_percentage = True
    value_axis.tick_labels.number_format = "0\%"  # This ensures the labels show as percentages


# ToDo: Was machen, wenn Daten transponiert sind
# ToDo: Labels und Axis formatting
def _create_bubble_chart(slide, df, headers, chart_core_message, header_cell_formats):
    data_selection_prompt = _create_bubble_chart_data_selection_prompt(headers, chart_core_message, "bubble chart",
                                                                       header_cell_formats)
    bubble_chart_data_structure = _query_openai(
        message=data_selection_prompt,
        response_model=BubbleChartDataStructure
    ) if not MOCK_AI_API_CALLS else (
        BubbleChartDataStructure(
            labels_column="Market",
            x_axis_column="Market share",
            y_axis_column="Market growth",
            x_axis_is_percentage=True,
            y_axis_is_percentage=True,
            x_axis_title="Market share (%)",
            y_axis_title="Market growth (%)",
            bubble_size_column="Market size",
            title="Market size in EUR"
        )
    )

    # Prepare data for the bubble chart
    pivot_df = df[[bubble_chart_data_structure.labels_column,
                   bubble_chart_data_structure.x_axis_column,
                   bubble_chart_data_structure.y_axis_column,
                   bubble_chart_data_structure.bubble_size_column]].copy()

    # Create chart data
    chart_data = BubbleChartData()

    # Convert all column headers to strings to avoid errors when matching with selected columns from openai prompt
    df.columns = df.columns.astype(str)

    # Add series and data points to the chart data
    for index, row in pivot_df.iterrows():
        category_label = row[bubble_chart_data_structure.labels_column]
        x_value = row[
            bubble_chart_data_structure.x_axis_column] if bubble_chart_data_structure.x_axis_is_percentage else row[
                                                                                                                    bubble_chart_data_structure.x_axis_column] * 100
        y_value = row[
            bubble_chart_data_structure.y_axis_column] if bubble_chart_data_structure.y_axis_is_percentage else row[
                                                                                                                    bubble_chart_data_structure.y_axis_column] * 100
        bubble_size = row[bubble_chart_data_structure.bubble_size_column]

        chart_data.add_series(category_label).add_data_point(x_value, y_value, bubble_size)

    # Chart creation
    diagram_placeholder = slide.placeholders[13]
    chart = diagram_placeholder.insert_chart(
        XL_CHART_TYPE.BUBBLE, chart_data
    ).chart

    # Title and labels
    slide.shapes.title.text = chart_core_message
    chart.has_title = True
    chart.chart_title.text_frame.text = bubble_chart_data_structure.title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = LABEL_COLOR
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    # Gridlines
    value_axis = chart.value_axis
    major_gridlines = value_axis.major_gridlines
    line = major_gridlines.format.line
    line.fill.solid()
    line.fill.fore_color.rgb = GRID_COLOR
    line.width = LINE_WIDTH

    # Style the x-axis
    category_axis = chart.category_axis
    axis_line = category_axis.format.line
    axis_line.fill.solid()
    category_axis.tick_labels.font.size = LABEL_FONT_SIZE
    category_axis.tick_labels.font.color.rgb = LABEL_COLOR
    category_axis.has_title = True
    category_axis.axis_title.text_frame.text = bubble_chart_data_structure.x_axis_title

    # Style the y-axis
    value_axis.tick_labels.font.size = LABEL_FONT_SIZE
    value_axis.tick_labels.font.color.rgb = LABEL_COLOR
    value_axis.format.line.fill.solid()
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = bubble_chart_data_structure.y_axis_title


# Prompts
def _create_chart_selection_prompt(df, chart_options, core_message, header_cell_formats):
    # Extract column names and descriptions
    columns = df.columns.tolist()

    # Summarize data overview (e.g., range or unique values for each column)
    data_overview = []
    for col in columns:
        if pd.api.types.is_numeric_dtype(df[col]):
            summary = f"Range: {df[col].min()} to {df[col].max()}"
        else:
            unique_vals = df[col].unique()
            unique_count = len(unique_vals)
            examples = ', '.join(map(str, unique_vals[:3]))  # Show up to 3 examples
            summary = f"Number of unique values: {unique_count} | Examples: {examples}"
        data_overview.append(f"- {col}: {summary}")
    data_overview_text = "\n".join(data_overview)

    # Include the first 5 rows of the DataFrame
    first_five_rows = df.head(5).to_string(index=False)

    # last 3 rows
    last_three_rows = df.tail(3).to_string(index=False)

    # Include header cell formats
    header_formats = "\n".join([f"- {col}: {fmt}" for col, fmt in header_cell_formats.items()])

    # Prepare chart options text
    chart_options_text = "\n".join(
        [
            f"*{option.value}*\n"
            f"   - Best suited for: {option.purpose}\n"
            f"   - Required data: {option.data_input}"
            for option in chart_options
        ]
    )

    # Construct the prompt
    prompt = f"""
I have a table with the following summary characteristics:

- **Column names:**
{columns}

- **Data overview:**
{data_overview_text}

- **Header cell formats:**
{header_formats}

- **First 5 rows of the data:**
{first_five_rows}

The chart should support the following message:
"{core_message}"

**Task:** Based on this summary of the table, please explain your reasoning and 
select the appropriate chart type from the following options:
{chart_options_text}

Consider the relationships and trends in the data to make your selection.


***The last 3 rows of the data:***
{last_three_rows}

Please determine if the last row contains the sum of all previous rows.
    """
    return prompt.strip()


def _create_bar_or_column_chart_data_selection_prompt(table_headers, chart_message, chart_type, header_cell_formats):
    return (
        f"You are provided with a table containing the following columns: {table_headers}.\n"
        f"Each column has a specific format, as described here: {header_cell_formats}.\n"
        f"The goal is to create a {chart_type} to support the following message:\n"
        f"'{chart_message}'.\n\n"
        f"Please identify:\n"
        f"1. Which column should be used as the category axis?\n"
        f"2. Which column should be used as the value axis?\n\n"
        f"The column names must match exactly the column names that were provided above."
        f"Additionally, provide a short descriptive name for the values, including a unit if applicable "
        f"(e.g., 'Living room size in square meters' or 'Vehicle sales in EUR'). "
        f"For currency units please always use the ISO currency code e.g. EUR instead of €"
    )


# No longer used
def _create_clustered_bar_or_column_chart_data_selection_prompt(table_headers, chart_message, chart_type,
                                                                header_cell_formats):
    return (
        f"You are provided with a table containing the following columns: {table_headers}.\n"
        f"Each column has a specific format, as described here: {header_cell_formats}.\n"
        f"The goal is to create a {chart_type} to support the following message:\n"
        f"'{chart_message}'.\n\n"
        f"Please identify:\n"
        f"1. Which column should be used as the categories?\n"
        f"2. Which column should be used as subcategories?\n"
        f"2. Which column should be used as values?\n\n"
        f"Additionally, provide a short descriptive name for the values, including a unit if applicable "
        f"(e.g., 'Living room size in square meters' or 'Vehicle sales in EUR'). For a stacked 100% chart the unit is always %."
        f"For currency units please always use the ISO currency code e.g. EUR instead of €"
    )


def _create_line_or_clustered_column_chart_data_selection_prompt(table_headers, chart_message, chart_type,
                                                                 header_cell_formats):
    return (
        f"You are provided with a table containing the following columns: {table_headers}.\n"
        f"Each column has a specific format, as described here: {header_cell_formats}.\n"
        f"The goal is to create a {chart_type} to support the following message:\n"
        f"'{chart_message}'.\n\n"
        f"Please identify:\n"
        f"1. Which column should be used as the categories?\n"
        f"2. Which columns should be used for the series data? Please list all columns. Do not include columns that "
        f"contain sums! \n"
        f"The column names must match exactly the column names that were provided above."
        f"3. Additionally, provide a short descriptive name for the values of the series data, including a unit if "
        f"applicable (e.g., 'Living room size in square meters' or 'Vehicle sales in EUR'). For 100% stacked column "
        f"chart, the unit is always %. For currency units please always use the ISO currency code e.g. EUR instead of €"
    )


def _create_bubble_chart_data_selection_prompt(table_headers, chart_message, chart_type, header_cell_formats):
    return (
        f"You are provided with a table containing the following columns: {table_headers}.\n"
        f"Each column has a specific format, as described here: {header_cell_formats}.\n"
        f"The goal is to create a {chart_type} to support the following message:\n"
        f"'{chart_message}'.\n\n"
        f"Please identify:\n"
        f"1. Which column should be used as the x-axis?\n"
        f"2. Which column should be used as the y-axis?\n"
        f"3. Which column should be used for the labels?\n"
        f"4. Which column should be used for the bubble size?\n"
        f"The column names must match exactly the column names that were provided above."
        f"5. For x- and y-axis and bubble size please provide a descriptive title if applicable (e.g., 'Living room "
        f"size in square meters' or 'Vehicle sales in EUR')."
        f"For currency units please always use the ISO currency code e.g. EUR instead of €"
        f"For percentages please use the % symbol"
        f"6. Additionally, provide a short descriptive title for the chart"
    )


def _create_pie_chart_data_selection_prompt(table_headers, chart_message, chart_type, header_cell_formats):
    return (
        f"You are provided with a table containing the following columns: {table_headers}.\n"
        f"Each column has a specific format, as described here: {header_cell_formats}.\n"
        f"The goal is to create a {chart_type} to support the following message:\n"
        f"'{chart_message}'.\n\n"
        f"Please identify:\n"
        f"1. Which column should be used as the x-axis?\n"
        f"2. Which column should be used as the y-axis?\n"
        f"The column names must match exactly the column names that were provided above."
        f"3. Please provide a descriptive title for the chart and include in the unit."
        f"The unit is always percent, please use the % symbol for that"
    )


def _create_pivot_data_prompt(df, core_message, chart_type):
    # Extract column names and descriptions
    columns = df.columns.tolist()

    # Summarize data overview (e.g., range or unique values for each column)
    data_overview = []
    for col in columns:
        if pd.api.types.is_numeric_dtype(df[col]):
            summary = f"Range: {df[col].min()} to {df[col].max()}"
        else:
            unique_vals = df[col].unique()
            summary = f"Examples: {', '.join(map(str, unique_vals[:3]))}"  # Show up to 3 examples
        data_overview.append(f"- {col}: {summary}")
    data_overview_text = "\n".join(data_overview)

    # Include the first 10 rows of the DataFrame
    first_ten_rows = df.head(10).to_string(index=False)

    # Construct the prompt
    prompt = f"""
    I have a table with the following summary characteristics:

    - **Column names:**
    {columns}

    - **Data overview:**
    {data_overview_text}

    - **First 10 rows of the data:**
    {first_ten_rows}

    I want to create a {chart_type} from it that supports the following message:
    "{core_message}"

    ### Questions:
    1. **Data structure analysis:**
        - Is the data currently in long format (i.e., multiple rows for each category)? 
   
    2. **LongFormFormat:**
       - Answer with "True" if the data is in long format.
    
    3. **Explain pivoting**
        - Assume the table needs to be pivoted in order to input it into the required chart type  
        - Explain which column should be the index, which column should be the new series and which column should be the values
         
    4. **Determine Columns for Pivoting:**
       - What column should be the **index**? (i.e., the x-axis values of the chart)
       - What column should be the **new columns**? (i.e., one series per column in the chart)
       - What column should be the **new values**? (i.e., the y-axis values of the chart)
    
    For the index, columns and values please answer with either a column name or a list of column names. 
    The column names must match exactly the column names that were provided above. Do not add quotation marks or anything else.
    """
    return prompt.strip()


# Data ingestion
def _extract_header_cell_formats(excel_bytes_content):
    workbook = load_workbook(excel_bytes_content)
    sheet = workbook.active

    # Create a dictionary mapping headers to the raw number formats of the second row
    headers_cellformatting_dict = {}
    for header_cell, data_cell in zip(sheet[1], sheet[2]):  # Row 1 for headers, Row 2 for formats
        headers_cellformatting_dict[header_cell.value] = data_cell.number_format

    # Output the headers and their raw formats
    return headers_cellformatting_dict


# Main function
def create_chart(excel_bytes_content, chart_core_message, uuid):
    presentation = Presentation(TEMPLATE_PATH)
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])

    df = pd.read_excel(excel_bytes_content)
    df_headers = df.columns.tolist()
    header_cell_formats = _extract_header_cell_formats(excel_bytes_content)

    has_more_than_two_headers = len(df_headers) > 2

    chart_selection_prompt = _create_chart_selection_prompt(
        df=df,
        chart_options=ChartType.get_all() if has_more_than_two_headers else ChartType.get_all_chart_types_with_two_columns_input(),
        core_message=chart_core_message,
        header_cell_formats=header_cell_formats)

    selected_chart_type = _query_openai(message=chart_selection_prompt, response_model=SelectedChartType)
    # selected_chart_type = SelectedChartType(
    #     reasonForSelectedChartType="some reason",
    #     chartType=ChartType.LINE,
    #     lastLineIncludesSum=True
    # )

    print(selected_chart_type.reasonForSelectedChartType)

    if selected_chart_type.lastLineIncludesSum:
        df = df.drop(df.index[-1])

    match selected_chart_type.chartType:
        case ChartType.COLUMN.value:
            _create_column_chart(
                slide=slide, df=df, headers=df_headers,
                chart_core_message=chart_core_message,
                header_cell_formats=header_cell_formats)
        case ChartType.COLUMN_CLUSTERED.value:
            _create_clustered_column_chart(
                slide=slide, df=df, headers=df_headers,
                chart_core_message=chart_core_message,
                header_cell_formats=header_cell_formats)
        case ChartType.COLUMN_STACKED.value:
            _create_stacked_column_chart(
                slide=slide, df=df, headers=df_headers,
                chart_core_message=chart_core_message,
                header_cell_formats=header_cell_formats)
        case ChartType.COLUMN_STACKED_100.value:
            _create_100_percent_stacked_column_chart(
                slide=slide, df=df, headers=df_headers,
                chart_core_message=chart_core_message,
                header_cell_formats=header_cell_formats)
        # case ChartType.BAR.value:
        #     _create_bar_chart(
        #         slide=slide, df=df, headers=df_headers,
        #         chart_core_message=chart_core_message,
        #         header_cell_formats=header_cell_formats)
        # case ChartType.BAR_CLUSTERED.value:
        #     _create_clustered_bar_chart(
        #         slide=slide, df=df, headers=df_headers,
        #         chart_core_message=chart_core_message,
        #         header_cell_formats=header_cell_formats)
        # case ChartType.BAR_STACKED.value:
        #     _create_stacked_bar_chart(
        #         slide=slide, df=df, headers=df_headers,
        #         chart_core_message=chart_core_message,
        #         header_cell_formats=header_cell_formats)
        # case ChartType.BAR_STACKED_100.value:
        #     _create_100_percent_stacked_bar_chart(
        #         slide=slide, df=df, headers=df_headers,
        #         chart_core_message=chart_core_message,
        #         header_cell_formats=header_cell_formats)
        case ChartType.PIE.value:
            _create_pie_chart(
                slide=slide, df=df, headers=df_headers,
                chart_core_message=chart_core_message,
                header_cell_formats=header_cell_formats)
        # case ChartType.DOUGHNUT.value:
        #     _create_doughnut_chart(
        #         slide=slide, df=df, headers=df_headers,
        #         chart_core_message=chart_core_message,
        #         header_cell_formats=header_cell_formats)
        case ChartType.LINE.value:
            _create_line_chart(
                slide=slide, df=df, headers=df_headers,
                chart_core_message=chart_core_message,
                header_cell_formats=header_cell_formats)
        # case ChartType.AREA_STACKED.value:
        #     _create_stacked_area_chart(
        #         slide=slide, df=df, headers=df_headers,
        #         chart_core_message=chart_core_message,
        #         header_cell_formats=header_cell_formats)
        case ChartType.AREA_STACKED_100.value:
            _create_100_percent_stacked_area_chart(
                slide=slide, df=df, headers=df_headers,
                chart_core_message=chart_core_message,
                header_cell_formats=header_cell_formats)
        case ChartType.BUBBLE.value:
            _create_bubble_chart(
                slide=slide, df=df, headers=df_headers,
                chart_core_message=chart_core_message,
                header_cell_formats=header_cell_formats)

    presentation_path = f"{uuid}_{selected_chart_type.chartType}_{datetime.datetime.now()}.pptx"
    presentation.save(presentation_path)
    return presentation_path

# create_chart(excel_bytes_content=DATA_TWO_TIME_SERIES,
#              chart_core_message="China will be the most important ice cream market by 2029",
#              uuid="1")
