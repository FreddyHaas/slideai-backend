# Column chart creators
from pptx.chart.data import CategoryChartData, BubbleChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Pt

from models import TwoColumnDataStructure, MultiColumnDataStructure, BubbleChartDataStructure, RoundingPrecision

AXIS_LABEL_COLOR = RGBColor(89, 89, 89)
DARK_GREEN = RGBColor(3, 90, 65)
DARK_GRAY = RGBColor(58, 58, 58)
MEDIUM_GRAY = RGBColor(116, 116, 116)
GRID_COLOR = RGBColor(217, 217, 217)
WHITE = RGBColor(255, 255, 255)
SIZE_12 = Pt(12)
SIZE_14 = Pt(14)
SIZE_18 = Pt(18)
LINE_WIDTH = Pt(0.4).emu


def create_column_chart(presentation, dataframe, chart_information: TwoColumnDataStructure, chart_core_message: str,
                        rounding_precision: RoundingPrecision):
    try:

        # Add slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])

        # Add chart to slide
        chart_data = CategoryChartData()
        chart_data.categories = dataframe[chart_information.category].tolist()
        chart_data.add_series(
            chart_information.value,
            dataframe[chart_information.value].tolist()
        )

        diagram_placeholder = slide.placeholders[13]
        chart = diagram_placeholder.insert_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data
        ).chart

        # Add action title
        slide.shapes.title.text = chart_core_message

        placeholder = slide.placeholders[1]
        _set_label(placeholder, chart_information, rounding_precision)

        # Remove Gridlines
        value_axis = chart.value_axis
        value_axis.visible = False
        value_axis.has_major_gridlines = False
        value_axis.has_minor_gridlines = False
        chart.has_title = False

        # Style the x-axis (category axis)
        category_axis = chart.category_axis

        # Format the x-axis line to match the major gridlines
        axis_line = category_axis.format.line
        axis_line.fill.solid()
        axis_line.fill.fore_color.rgb = GRID_COLOR  # Same color as gridlines
        axis_line.width = LINE_WIDTH  # Same width as gridlines

        category_labels = category_axis.tick_labels.font
        category_labels.color.rgb = DARK_GRAY
        category_labels.bold = True

        no_of_categories = len(dataframe)
        category_labels.size = Pt(14) if no_of_categories < 11 else Pt(12)

        # Add data labels
        for series in chart.series:
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.show_value = True
            data_labels.font.size = Pt(16) if no_of_categories < 11 else Pt(12)
            data_labels.font.bold = True
            data_labels.font.color.rgb = DARK_GREEN
            data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END  # Position labels outside the bars
            data_labels.number_format = _resolve_number_format(rounding_precision=rounding_precision)

            for point in series.points:
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = DARK_GREEN

    except Exception as exception:
        _delete_last_slide(presentation)
        print(str(exception))


def create_clustered_column_chart(presentation, dataframe, chart_information, chart_core_message,
                                  rounding_precision: RoundingPrecision):
    try:
        # Add slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])

        # Chart creation
        chart_data = CategoryChartData()
        chart_data.categories = dataframe[chart_information.category].tolist()

        no_of_entries = len(chart_information.series) * len(chart_data.categories)

        for column in chart_information.series:
            chart_data.add_series(column, dataframe[column].tolist())

        diagram_placeholder = slide.placeholders[13]
        chart = diagram_placeholder.insert_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data
        ).chart

        # Title and labels
        slide.shapes.title.text = chart_core_message
        chart.has_title = False

        placeholder = slide.placeholders[1]
        _set_label(placeholder, chart_information, rounding_precision)

        # Chart legend
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM  # Bottom of the chart (default)
        chart.legend.font.size = SIZE_12

        # Gridlines
        value_axis = chart.value_axis
        major_gridlines = value_axis.major_gridlines
        line = major_gridlines.format.line
        line.fill.solid()
        line.fill.fore_color.rgb = GRID_COLOR
        line.width = LINE_WIDTH

        # Style the x-axis (category axis)
        category_axis = chart.category_axis

        # Format the x-axis line to match the major gridlines
        axis_line = category_axis.format.line
        axis_line.fill.solid()
        axis_line.fill.fore_color.rgb = GRID_COLOR  # Same color as gridlines
        axis_line.width = LINE_WIDTH  # Same width as gridlines

        category_axis.tick_labels.font.size = SIZE_12
        category_axis.tick_labels.font.color.rgb = DARK_GRAY
        category_axis.tick_labels.font.bold = True

        # Style the y-axis (value axis) but hide its line
        value_axis.tick_labels.font.size = SIZE_12
        value_axis.tick_labels.font.color.rgb = AXIS_LABEL_COLOR
        value_axis.format.line.fill.solid()
        value_axis.format.line.fill.background()  # Hide the axis line

        if no_of_entries < 20:

            value_axis.visible = False
            value_axis.has_minor_gridlines = False
            value_axis.has_major_gridlines = False

            # Add data labels
            for series in chart.series:
                series.has_data_labels = True
                data_labels = series.data_labels
                data_labels.show_value = True
                data_labels.font.size = Pt(12)
                data_labels.font.bold = True
                data_labels.font.color.rgb = DARK_GREEN
                data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                data_labels.number_format = _resolve_number_format(rounding_precision=rounding_precision)

    except Exception as exception:
        _delete_last_slide(presentation)
        print(str(exception))


def create_stacked_column_chart(presentation, dataframe, chart_information, chart_core_message,
                                rounding_precision: RoundingPrecision):
    try:
        # Add slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])

        no_of_series = len(chart_information.series)

        # Chart creation
        chart_data = CategoryChartData()
        chart_data.categories = dataframe[chart_information.category].tolist()

        for column in chart_information.series:
            chart_data.add_series(str(column), dataframe[column].tolist())

        diagram_placeholder = slide.placeholders[13]
        chart = diagram_placeholder.insert_chart(
            XL_CHART_TYPE.COLUMN_STACKED, chart_data  # Set chart type to stacked column
        ).chart

        chart.plots[0].gap_width = 100

        # Title and labels
        slide.shapes.title.text = chart_core_message

        placeholder = slide.placeholders[1]
        _set_label(placeholder, chart_information, rounding_precision)

        # Chart legend
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM  # Bottom of the chart (default)
        chart.legend.font.size = SIZE_12

        # Gridlines
        value_axis = chart.value_axis
        value_axis.visible = False
        value_axis.has_major_gridlines = False
        value_axis.has_minor_gridlines = False

        # Style the x-axis (category axis)
        category_axis = chart.category_axis

        # Format the x-axis line to match the major gridlines
        axis_line = category_axis.format.line
        axis_line.fill.solid()
        axis_line.fill.fore_color.rgb = GRID_COLOR  # Same color as gridlines
        axis_line.width = LINE_WIDTH  # Same width as gridlines

        category_axis.tick_labels.font.size = SIZE_14
        category_axis.tick_labels.font.color.rgb = DARK_GRAY
        category_axis.tick_labels.font.bold = True

        for series in chart.series:
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.show_value = True
            data_labels.font.size = Pt(14) if no_of_series < 4 else Pt(12)
            data_labels.font.bold = True
            data_labels.font.color.rgb = WHITE
            data_labels.position = XL_DATA_LABEL_POSITION.CENTER
            data_labels.number_format = _resolve_number_format(rounding_precision=rounding_precision)

    except Exception as exception:
        _delete_last_slide(presentation)
        print(str(exception))


def create_100_percent_stacked_column_chart(presentation, dataframe, chart_information, chart_core_message):
    try:
        # Add slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])

        no_of_series = len(chart_information.series)

        # Chart creation
        chart_data = CategoryChartData()
        chart_data.categories = dataframe[chart_information.category].tolist()

        for column in chart_information.series:
            chart_data.add_series(str(column), dataframe[column].tolist())

        diagram_placeholder = slide.placeholders[13]
        chart = diagram_placeholder.insert_chart(
            XL_CHART_TYPE.COLUMN_STACKED_100, chart_data  # Set chart type to stacked column
        ).chart

        chart.plots[0].gap_width = 100

        # Title and labels
        slide.shapes.title.text = chart_core_message

        # Chart legend
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM  # Bottom of the chart (default)
        chart.legend.font.size = SIZE_12

        # Gridlines
        value_axis = chart.value_axis
        value_axis.visible = False
        value_axis.has_major_gridlines = False
        value_axis.has_minor_gridlines = False

        # Style the x-axis (category axis)
        category_axis = chart.category_axis

        # Format the x-axis line to match the major gridlines
        axis_line = category_axis.format.line
        axis_line.fill.solid()
        axis_line.fill.fore_color.rgb = GRID_COLOR  # Same color as gridlines
        axis_line.width = LINE_WIDTH  # Same width as gridlines

        category_axis.tick_labels.font.size = SIZE_14
        category_axis.tick_labels.font.color.rgb = DARK_GRAY
        category_axis.tick_labels.font.bold = True

        for series in chart.series:
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.show_value = True
            data_labels.font.size = Pt(14) if no_of_series < 4 else Pt(12)
            data_labels.font.bold = True
            data_labels.font.color.rgb = WHITE
            data_labels.position = XL_DATA_LABEL_POSITION.CENTER
            data_labels.number_format = '0'

    except Exception as exception:
        _delete_last_slide(presentation)
        print(str(exception))


# Bar chart creators
def create_bar_chart(presentation, dataframe, chart_information, chart_core_message,
                     rounding_precision: RoundingPrecision):
    try:

        # Add slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])

        # Add chart to slide
        chart_data = CategoryChartData()
        chart_data.categories = dataframe[chart_information.category].tolist()
        chart_data.add_series(
            chart_information.value,
            dataframe[chart_information.value].tolist()
        )

        diagram_placeholder = slide.placeholders[13]
        chart = diagram_placeholder.insert_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, chart_data
        ).chart

        # Add action title
        slide.shapes.title.text = chart_core_message

        slide.placeholders[1].text = chart_information.axis_label

        placeholder = slide.placeholders[1]
        _set_label(placeholder, chart_information, rounding_precision)

        # Remove minor gridlines and titles
        value_axis = chart.value_axis
        value_axis.visible = False
        value_axis.has_minor_gridlines = False
        chart.has_title = False

        # Style the major gridlines
        value_axis.has_major_gridlines = True
        major_gridlines = value_axis.major_gridlines.format.line
        major_gridlines.fill.solid()
        major_gridlines.fill.fore_color.rgb = GRID_COLOR
        major_gridlines.width = LINE_WIDTH

        # Style the x-axis (category axis)
        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False

        # Format the x-axis line to match the major gridlines
        axis_line = category_axis.format.line
        axis_line.fill.solid()
        axis_line.fill.fore_color.rgb = GRID_COLOR  # Same color as gridlines
        axis_line.width = LINE_WIDTH  # Same width as gridlines

        category_labels = category_axis.tick_labels.font
        category_labels.color.rgb = DARK_GRAY

        no_of_categories = len(dataframe)
        category_labels.size = Pt(14) if no_of_categories < 16 else Pt(10)
        category_labels.bold = True if no_of_categories < 11 else False

        # Add data labels
        for series in chart.series:
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.show_value = True
            data_labels.font.size = Pt(16) if no_of_categories < 11 else Pt(14) if no_of_categories < 16 else Pt(12)
            data_labels.font.bold = True
            data_labels.font.color.rgb = DARK_GREEN
            data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END  # Position labels outside the bars
            data_labels.number_format = _resolve_number_format(rounding_precision=rounding_precision)

            for point in series.points:
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = DARK_GREEN

    except Exception as exception:
        _delete_last_slide(presentation)
        print(str(exception))


def create_clustered_bar_chart(presentation, dataframe, chart_information, chart_core_message,
                               rounding_precision: RoundingPrecision):
    try:
        # Add slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])

        # Chart creation
        chart_data = CategoryChartData()
        chart_data.categories = dataframe[chart_information.category].tolist()

        for column in chart_information.series:
            chart_data.add_series(str(column), dataframe[column].tolist())

        diagram_placeholder = slide.placeholders[13]
        chart = diagram_placeholder.insert_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, chart_data
        ).chart

        # Title and labels
        slide.shapes.title.text = chart_core_message

        placeholder = slide.placeholders[1]
        _set_label(placeholder, chart_information, rounding_precision)

        # Chart legend
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = SIZE_12

        # Gridlines
        value_axis = chart.value_axis
        major_gridlines = value_axis.major_gridlines
        line = major_gridlines.format.line
        line.fill.solid()
        line.fill.fore_color.rgb = GRID_COLOR
        line.width = LINE_WIDTH

        # Style the y-axis (category axis)
        category_axis = chart.category_axis

        # Format the y-axis line to match the major gridlines
        axis_line = category_axis.format.line
        axis_line.fill.solid()
        axis_line.fill.fore_color.rgb = GRID_COLOR
        axis_line.width = LINE_WIDTH

        category_axis.tick_labels.font.size = SIZE_12
        category_axis.tick_labels.font.color.rgb = DARK_GRAY
        category_axis.tick_labels.font.bold = True

        # Style the x-axis (value axis)
        value_axis.tick_labels.font.size = SIZE_12
        value_axis.tick_labels.font.color.rgb = AXIS_LABEL_COLOR
        value_axis.format.line.fill.solid()
        value_axis.format.line.fill.background()

        no_of_entries = len(chart_information.series) * len(chart_data.categories)

        if no_of_entries < 11:

            value_axis.visible = False
            value_axis.has_major_gridlines = False
            value_axis.has_minor_gridlines = False

            # Add data labels
            for series in chart.series:
                series.has_data_labels = True
                data_labels = series.data_labels
                data_labels.show_value = True
                data_labels.font.size = Pt(12)
                data_labels.font.bold = True
                data_labels.font.color.rgb = DARK_GREEN
                data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                data_labels.number_format = _resolve_number_format(rounding_precision=rounding_precision)

    except Exception as exception:
        _delete_last_slide(presentation)
        print(str(exception))


def create_stacked_bar_chart(presentation, dataframe, chart_information, chart_core_message,
                             rounding_precision: RoundingPrecision):
    try:
        # Add slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])

        no_of_series = len(chart_information.series)

        # Chart creation
        chart_data = CategoryChartData()
        chart_data.categories = dataframe[chart_information.category].tolist()

        for column in chart_information.series:
            chart_data.add_series(str(column), dataframe[column].tolist())

        diagram_placeholder = slide.placeholders[13]
        chart = diagram_placeholder.insert_chart(
            XL_CHART_TYPE.BAR_STACKED, chart_data  # Set chart type to stacked column
        ).chart

        chart.plots[0].gap_width = 100

        # Title and labels
        slide.shapes.title.text = chart_core_message

        placeholder = slide.placeholders[1]
        _set_label(placeholder, chart_information, rounding_precision)

        # Chart legend
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM  # Bottom of the chart (default)
        chart.legend.font.size = SIZE_12

        # Gridlines
        value_axis = chart.value_axis
        value_axis.visible = False
        value_axis.has_major_gridlines = False
        value_axis.has_minor_gridlines = False

        # Style the x-axis (category axis)
        category_axis = chart.category_axis

        # Format the x-axis line to match the major gridlines
        axis_line = category_axis.format.line
        axis_line.fill.solid()
        axis_line.fill.fore_color.rgb = GRID_COLOR  # Same color as gridlines
        axis_line.width = LINE_WIDTH  # Same width as gridlines

        category_axis.tick_labels.font.size = SIZE_14
        category_axis.tick_labels.font.color.rgb = DARK_GRAY
        category_axis.tick_labels.font.bold = True

        for series in chart.series:
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.show_value = True
            data_labels.font.size = Pt(14) if no_of_series < 4 else Pt(12)
            data_labels.font.bold = True
            data_labels.font.color.rgb = WHITE
            data_labels.position = XL_DATA_LABEL_POSITION.CENTER
            data_labels.number_format = _resolve_number_format(rounding_precision=rounding_precision)

    except Exception as exception:
        _delete_last_slide(presentation)
        print(str(exception))


def create_100_percent_stacked_bar_chart(presentation, dataframe, chart_information: MultiColumnDataStructure,
                                         chart_core_message: str):
    try:
        # Add slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])

        pivot_df = dataframe.sort_values(by=dataframe.columns[-1], ascending=True)

        # Normalize the values to percentages for 100% stacked bar chart
        pivot_df_percentage = pivot_df.copy()
        pivot_df_percentage[chart_information.series] = pivot_df[
                                                            chart_information.series].div(
            pivot_df[chart_information.series].sum(axis=1), axis=0
        ) * 100  # Convert to percentage

        # Prepare the data for the chart
        chart_data = CategoryChartData()
        chart_data.categories = pivot_df_percentage[chart_information.category].tolist()

        # Add each subcategory as a series for stacking (now with percentage values)
        for column in chart_information.series:
            chart_data.add_series(column, pivot_df_percentage[column].tolist())

        # Insert a 100% stacked bar chart
        diagram_placeholder = slide.placeholders[13]
        chart = diagram_placeholder.insert_chart(
            XL_CHART_TYPE.BAR_STACKED_100, chart_data
        ).chart

        # Title and labels
        slide.shapes.title.text = chart_core_message
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_information.title
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = AXIS_LABEL_COLOR
        chart.chart_title.text_frame.paragraphs[0].font.bold = False

        # Chart legend
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = SIZE_12

        # Gridlines
        value_axis = chart.value_axis
        major_gridlines = value_axis.major_gridlines
        line = major_gridlines.format.line
        line.fill.solid()
        line.fill.fore_color.rgb = GRID_COLOR
        line.width = LINE_WIDTH

        # Style the y-axis (category axis)
        category_axis = chart.category_axis
        axis_line = category_axis.format.line
        axis_line.fill.solid()
        axis_line.fill.fore_color.rgb = GRID_COLOR
        axis_line.width = LINE_WIDTH

        category_axis.tick_labels.font.size = SIZE_12
        category_axis.tick_labels.font.color.rgb = AXIS_LABEL_COLOR

        # Style the x-axis (value axis)
        value_axis.tick_labels.font.size = SIZE_12
        value_axis.tick_labels.font.color.rgb = AXIS_LABEL_COLOR
        value_axis.format.line.fill.solid()
        value_axis.format.line.fill.background()  # Hide the axis line

        # Adjust gap width for stacking aesthetics
        chart.plots[0].gap_width = 50  # Adjust as needed for aesthetics
    except Exception as exception:
        _delete_last_slide(presentation)
        print(str(exception))


# Pie chart creators
def create_pie_chart(presentation, dataframe, chart_information: TwoColumnDataStructure, chart_core_message):
    try:
        # Add slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])

        # Convert all column headers to strings to avoid errors when matching with selected columns from openai prompt
        dataframe.columns = dataframe.columns.astype(str)

        # Create a CategoryChartData object
        chart_data = CategoryChartData()
        chart_data.categories = dataframe[chart_information.category].tolist()
        chart_data.add_series('Percentage', dataframe[chart_information.value].tolist())

        # Create the pie chart
        diagram_placeholder = slide.placeholders[13]
        chart = diagram_placeholder.insert_chart(XL_CHART_TYPE.PIE, chart_data).chart

        # Title and labels
        slide.shapes.title.text = chart_core_message
        chart.has_title = False

        # Chart legend
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(10)  # You can customize this value

        # Adding data labels to the pie chart
        plot = chart.plots[0]  # Pie chart typically has one plot
        plot.has_data_labels = True  # Enable data labels
        data_labels = plot.data_labels
        data_labels.show_value = True  # Show values
        data_labels.number_format = "0%"  # Customize the format, e.g., "Percentage" for percent values
        data_labels.font.size = Pt(16)  # Customize the font size of the labels
        data_labels.font.bold = True
        data_labels.font.color.rgb = WHITE

    except Exception as exception:
        _delete_last_slide(presentation)
        print(str(exception))


def create_doughnut_chart(presentation, dataframe, chart_information: TwoColumnDataStructure, chart_core_message):
    try:
        # Add slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])

        # Convert all column headers to strings to avoid errors when matching with selected columns from openai prompt
        dataframe.columns = dataframe.columns.astype(str)

        # Create a CategoryChartData object
        chart_data = CategoryChartData()
        chart_data.categories = dataframe[chart_information.category].tolist()
        chart_data.add_series('Percentage', dataframe[chart_information.value].tolist())

        # Create the pie chart
        diagram_placeholder = slide.placeholders[13]
        chart = diagram_placeholder.insert_chart(XL_CHART_TYPE.DOUGHNUT, chart_data).chart

        # Title and labels
        slide.shapes.title.text = chart_core_message
        chart.has_title = False

        # Chart legend
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(10)  # You can customize this value

        # Adding data labels to the pie chart
        plot = chart.plots[0]  # Pie chart typically has one plot
        plot.has_data_labels = True  # Enable data labels
        data_labels = plot.data_labels
        data_labels.show_value = True  # Show values
        data_labels.number_format = "0%"  # Customize the format, e.g., "Percentage" for percent values
        data_labels.font.size = Pt(16)  # Customize the font size of the labels
        data_labels.font.bold = True
        data_labels.font.color.rgb = WHITE

    except Exception as exception:
        _delete_last_slide(presentation)
        print(str(exception))


# Time series data
def create_line_chart(presentation, dataframe, chart_information, chart_core_message):
    try:
        # Add slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])

        # Chart creation
        chart_data = CategoryChartData()
        chart_data.categories = dataframe[chart_information.category].tolist()

        for column in chart_information.series:
            chart_data.add_series(column, dataframe[column].tolist())

        diagram_placeholder = slide.placeholders[13]
        chart = diagram_placeholder.insert_chart(
            XL_CHART_TYPE.LINE, chart_data
        ).chart

        # Title and labels
        slide.shapes.title.text = chart_core_message
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_information.axis_label
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = AXIS_LABEL_COLOR
        chart.chart_title.text_frame.paragraphs[0].font.bold = False

        # Chart legend
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = SIZE_12

        # Gridlines
        value_axis = chart.value_axis
        major_gridlines = value_axis.major_gridlines
        line = major_gridlines.format.line
        line.fill.solid()
        line.fill.fore_color.rgb = GRID_COLOR
        line.width = LINE_WIDTH

        # Style the y-axis (category axis)
        category_axis = chart.category_axis

        # Format the y-axis line to match the major gridlines
        axis_line = category_axis.format.line
        axis_line.fill.solid()
        axis_line.fill.fore_color.rgb = GRID_COLOR
        axis_line.width = LINE_WIDTH

        category_axis.tick_labels.font.size = SIZE_12
        category_axis.tick_labels.font.color.rgb = AXIS_LABEL_COLOR

        # Style the x-axis (value axis)
        value_axis.tick_labels.font.size = SIZE_12
        value_axis.tick_labels.font.color.rgb = AXIS_LABEL_COLOR
        value_axis.format.line.fill.solid()
        value_axis.format.line.fill.background()  # Hide the axis line
    except Exception as exception:
        _delete_last_slide(presentation)
        print(str(exception))


def create_bubble_chart(presentation, dataframe, chart_information: BubbleChartDataStructure, chart_core_message):
    try:
        # Add slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[2])

        # Create chart data
        chart_data = BubbleChartData()

        # Add series and data points to the chart data
        for index, row in dataframe.iterrows():
            category_label = row[chart_information.labels_column]
            x_value = row[chart_information.x_axis_column] \
                if chart_information.x_axis_is_percentage else row[chart_information.x_axis_column] * 100
            y_value = row[
                chart_information.y_axis_column] \
                if chart_information.y_axis_is_percentage else row[chart_information.y_axis_column] * 100
            bubble_size = row[chart_information.bubble_size_column]

            chart_data.add_series(category_label).add_data_point(x_value, y_value, bubble_size)

        # Chart creation
        diagram_placeholder = slide.placeholders[13]
        chart = diagram_placeholder.insert_chart(
            XL_CHART_TYPE.BUBBLE, chart_data
        ).chart

        # Title and labels
        slide.shapes.title.text = chart_core_message
        chart.has_title = False

        # Gridlines
        value_axis = chart.value_axis
        major_gridlines = value_axis.major_gridlines
        line = major_gridlines.format.line
        line.fill.solid()
        line.fill.fore_color.rgb = GRID_COLOR
        line.width = LINE_WIDTH

        # Style the x-axis
        category_axis = chart.category_axis
        axis_line = category_axis.format.line
        axis_line.fill.solid()
        category_axis.tick_labels.font.size = SIZE_12
        category_axis.tick_labels.font.color.rgb = AXIS_LABEL_COLOR
        category_axis.has_title = False

        # Style the y-axis
        value_axis.tick_labels.font.size = SIZE_12
        value_axis.tick_labels.font.color.rgb = AXIS_LABEL_COLOR
        value_axis.format.line.fill.solid()
        chart.value_axis.has_title = False
    except Exception as exception:
        _delete_last_slide(presentation)
        print(str(exception))


def _delete_last_slide(presentation):
    slide_id_list = presentation.slides._sldIdLst
    slide_id_list.remove(slide_id_list[-1])


def _resolve_number_format(rounding_precision: RoundingPrecision) -> str:
    order_of_magnitude = rounding_precision.order_of_magnitude
    decimal_place = rounding_precision.decimal_place
    number_format = '0'

    if decimal_place != 0:
        number_format = number_format + '.0'

    if order_of_magnitude < 3:
        return number_format

    if order_of_magnitude < 6:
        return number_format + ','

    if order_of_magnitude < 9:
        return number_format + ',,'

    return number_format + ',,,'


def _resolve_unit_label(unit: str, order_of_magnitude: int) -> str:
    # Determine magnitude based on order_of_magnitude
    if 3 <= order_of_magnitude <= 5:
        magnitude = "k"
    elif 6 <= order_of_magnitude <= 8:
        magnitude = "mn"
    elif order_of_magnitude > 8:
        magnitude = "bn"
    else:
        magnitude = ""

    unit_is_none = unit.lower() == "none" or unit.lower() == '"none"'

    # Return appropriate label based on conditions
    if unit_is_none and not magnitude:
        return ""
    elif unit_is_none and magnitude:
        return f"in {magnitude}"
    elif not unit_is_none and not magnitude:
        return f"in {unit}"
    else:
        return f"in {magnitude} {unit}"


def _set_label(placeholder, chart_information, rounding_precision):

    placeholder.text = chart_information.axis_label
    text_frame = placeholder.text_frame
    unit_text = text_frame.add_paragraph()
    unit_text.text = _resolve_unit_label(
        chart_information.axis_unit,
        rounding_precision.order_of_magnitude)
    unit_text.font.color.rgb = MEDIUM_GRAY
    unit_text.font.bold = False
    unit_text.space_before = Pt(0)
    unit_text.space_after = Pt(0)
